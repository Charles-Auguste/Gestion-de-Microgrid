# -*- coding: utf-8 -*-
"""
Created on Wed May 12 05:48:03 2021

@author: B57876
"""

# Create a .ppt to summarize results of a run of the microgrid(s)

# IMPORT
# standard
import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
import datetime
from datetime import timedelta
import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
# perso
from calc_output_metrics import check_if_unique_list
from plot import plot_mg_load_during_coord_method, plot_all_teams_mg_load_last_iter, \
            plot_per_actor_load_last_iter, plot_all_teams_cost_auton_tradeoff_last_iter, \
            plot_all_teams_score_traj

# Cf. https://python-pptx.readthedocs.io/en/latest/
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
def create_summary_of_run_ppt(result_dir: str, date_of_run: datetime.datetime,
                              idx_run: int, optim_period: pd.date_range,
                              coord_method: str, regions_map_file: str,
                              pv_prof: dict, load_profiles: dict,
                              microgrid_prof: dict, microgrid_pmax: dict,
                              cost_autonomy_tradeoff: dict, team_scores: dict,
                              best_teams_per_region: dict, podium_france_file: str,
                              teams_france_classif: dict, type_of_score: str, 
                              scores_traj: dict):
    """ 
    Create a powerpoint to summarize the results of a given run of the microgrid 
    serious game

    :param result_dir: full path to the directory where results are stored
    :param date_of_run: date of the current run 
    :param idx_run: index of current run (to get different "stages" of the tests)
    :param optim_period: optim. discrete-time period
    :param coord_method: name of the coordination method used in the microgrid
    :param regions_map_file: full path to the (French) regions map file
    :param pv_prof: PV prod. profile
    :param load_profiles: dictionary with keys 1. Industrial Cons. scenario; 
    2. Data Center scenario; 3. PV scenario; 4. EV scenario; 5. microgrid team name; 
    6. Iteration; 7. actor type (N.B. charging_station_1, charging_station_2... 
    if multiple charging stations in this microgrid) and values the associated
    load profile (kW)
    :param microgrid_prof: dict. with keys 1. Industrial Cons. scenario; 
    2. Data Center scenario; 3. PV scenario; 4. EV scenario; 5. microgrid team name; 
    6. Iteration and value the associated aggreg. microgrid load profile
    :param microgrid_pmax: dict. with keys 1. Industrial Cons. scenario; 
    2. Data Center scenario; 3. PV scenario; 4. EV scenario; 5. microgrid team name; 
    6. Iteration and value the associated microgrid pmax
    :param cost_autonomy_tradeoff: dict. with keys 1. the team names; 2. PV region
    names and values the associated (cost, autonomy score) aggreg. over the set 
    of other scenarios
    :param team_scores: dict. with keys 1. team name; 2. region name and values 
    the associated score of the current run
    :param best_teams_per_region: dict. with keys the names of the region simulated
    in current run and values the associated list of best team(s)
    :param podium_france_file: full path to the image of the "podium France" file
    :param teams_france_classif: dict. with keys the different rankings (from 1
    to the number of teams) and values the couple (list of the name(s) of the team(s)),
    score, which may be a number or a vector)
    :param type_of_score: the type of score used for classification (unique 
    number, or a vector)    
    :param scores_traj: dict. with keys 1. team name ; 2. run dates and associated
    values the score
    """
    
    img_format = "png"
    region_coord_dyn_plot = None # if None the first region will be used
    img_slide_layout_idx = 1 # layout used for the slides with images
#    table_slide_layout_idx = 1 # and the one for slides with tables
    left_empty_space = 2 # in percentage of slide width
    bottom_empty_space = 5 # idem with height
    font_name = "Calibri"
    font_size = 30
    font_bold = True
    font_italic = False
    text_vertical_align = "middle" # "top", "middle", "bottom"
    
    date_format = "%Y-%m-%d %H:%M"
    # for the table not to be too width
#    short_region_names = {"grand_nord": "gd N", "grand_est": "gd E",
#                          "grand_rhone": "gd Rhone", "bretagne": "bret.",
#                          "grand_ouest": "gd O", "grand_sud_ouest": "gd SO",
#                          "grande_ardeche": "gde ard.", "grand_sud_est": "gd SE"} 
    
    # create an empty presentation
    prs = Presentation()
        
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Summary of microgrid serious game run %i" % idx_run
    subtitle.text = ("Generated on {:%s}" % date_format).format(date_of_run)
    
    # DETAILED results
    # get team and region names
    team_names = list(team_scores)
    n_teams = len(team_names)
    regions = [list(team_scores[team]) for team in team_names]
    
    # check that all teams have the same regions scenarios
    if not check_if_unique_list(regions):
        print("Different teams does not have the same region names... -> STOP")
        sys.exit(1)
    else:
        regions=regions[0]

    # set region for coord dyn plot if not provided in parameter
    if region_coord_dyn_plot == None:
        region_coord_dyn_plot = regions[0]
    print("Region used for coordination dyn. plot: %s" % region_coord_dyn_plot)
    
    # TODO 2021-5-17 exclude teams with bug
    # slide to list these teams
    
    # 1) 1 slide per team with per iteration total microgrid load
    for i_team in range(n_teams):
        team = team_names[i_team]
        slide, shapes, title_shape = \
            init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                      "Team %s load DURING %s \n (PV) region: %s" \
                                      % (team, coord_method, region_coord_dyn_plot),
                                      font_name, font_size, font_bold, font_italic,
                                      text_vertical_align)
    
        # plot and save 
        current_dyn_mg_load_file = os.path.join(result_dir,
                                                "mg_load_during_dyn_team_%s_%s.%s" \
                                                % (team, date_of_run.strftime("%Y-%m-%d_%H%M"),
                                                   img_format))
        plot_mg_load_during_coord_method(microgrid_prof, region_coord_dyn_plot,
                                         team, current_dyn_mg_load_file.split(".")[0],
                                         optim_period)
        
        # open as Pillow Image
        dyn_mg_load_img = Image.open(current_dyn_mg_load_file)
        
        # add image
        add_img_to_slide(slide, dyn_mg_load_img, current_dyn_mg_load_file,
                         ((1-2*left_empty_space/100)*prs.slide_width,
                          (1-2*bottom_empty_space/100)*prs.slide_height \
                                                            -title_shape.height),
                          title_shape.height, left_empty_space/100*prs.slide_width,
                          bottom_empty_space/100*prs.slide_height)
        
        # suppress unused text placeholder (of index 1, 0 is used for the title)
        sp = shapes[1].element
        sp.getparent().remove(sp)
        
    # 1) 1 slide per team with last iteration per-actor load
    # TODO add possibility to plot best iter
    for i_team in range(n_teams):
        team = team_names[i_team]
        slide, shapes, title_shape = \
            init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                      "Team %s ACTORS' load AT THE END of %s \n (PV) region: %s" \
                                      % (team, coord_method, region_coord_dyn_plot),
                                      font_name, font_size, font_bold, font_italic,
                                      text_vertical_align)
    
        # plot and save 
        per_actor_load_file = os.path.join(result_dir,
                                                "per_actor_load_last_iter_team_%s_%s.%s" \
                                                % (team, date_of_run.strftime("%Y-%m-%d_%H%M"),
                                                   img_format))
        plot_per_actor_load_last_iter(load_profiles, pv_prof[region_coord_dyn_plot], region_coord_dyn_plot,
                                      team, per_actor_load_file.split(".")[0],
                                      optim_period)
        # open as Pillow Image
        per_actor_load_img = Image.open(per_actor_load_file)
        
        # add image
        add_img_to_slide(slide, per_actor_load_img, per_actor_load_file,
                         ((1-2*left_empty_space/100)*prs.slide_width,
                          (1-2*bottom_empty_space/100)*prs.slide_height \
                                                            -title_shape.height),
                          title_shape.height, left_empty_space/100*prs.slide_width,
                          bottom_empty_space/100*prs.slide_height)
        
        # suppress unused text placeholder (of index 1, 0 is used for the title)
        sp = shapes[1].element
        sp.getparent().remove(sp)
    

    # 1 slide with all teams microgrid load at the last iteration
    # TODO add possibility to plot best iter    
    slide, shapes, title_shape = \
        init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                  "All teams load AT THE END of %s \n (PV) region: %s" \
                                  % (coord_method, region_coord_dyn_plot),
                                  font_name, font_size, font_bold, font_italic,
                                  text_vertical_align)

    # plot and save 
    all_teams_mg_load_file = os.path.join(result_dir, "all_teams_mg_load_last_iter_%s.%s" \
                                           % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                              img_format))
    plot_all_teams_mg_load_last_iter(microgrid_prof, microgrid_pmax, pv_prof[region_coord_dyn_plot],
                                     region_coord_dyn_plot,
                                     all_teams_mg_load_file.split(".")[0],
                                     optim_period)
    
    # open as Pillow Image
    all_teams_mg_load_img = Image.open(all_teams_mg_load_file)
    
    # add image
    add_img_to_slide(slide, all_teams_mg_load_img, all_teams_mg_load_file,
                     ((1-2*left_empty_space/100)*prs.slide_width,
                      (1-2*bottom_empty_space/100)*prs.slide_height \
                                                        -title_shape.height),
                      title_shape.height, left_empty_space/100*prs.slide_width,
                      bottom_empty_space/100*prs.slide_height)
    
    # suppress unused text placeholder (of index 1, 0 is used for the title)
    sp = shapes[1].element
    sp.getparent().remove(sp)

    # 1 slide with scatter with 1 (eur, autonomy) point per team*PV region
    slide, shapes, title_shape = \
        init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                  "All teams (cost, autonomy) tradeoff with %s" \
                                  % coord_method, font_name, font_size, font_bold, 
                                  font_italic, text_vertical_align)

    # plot and save 
    all_teams_cost_auton_tradeoff_file = os.path.join(result_dir,
                                                      "all_teams_cost_auton_tradeoff_last_iter_%s.%s" \
                                           % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                              img_format))

    plot_all_teams_cost_auton_tradeoff_last_iter(cost_autonomy_tradeoff,
                                                 all_teams_cost_auton_tradeoff_file.split(".")[0])
    
    # open as Pillow Image
    all_teams_cost_auton_tradeoff_img = Image.open(all_teams_cost_auton_tradeoff_file)
    
    # add image
    add_img_to_slide(slide, all_teams_cost_auton_tradeoff_img,
                     all_teams_cost_auton_tradeoff_file,
                     ((1-2*left_empty_space/100)*prs.slide_width,
                      (1-2*bottom_empty_space/100)*prs.slide_height \
                                                        -title_shape.height),
                      title_shape.height, left_empty_space/100*prs.slide_width,
                      bottom_empty_space/100*prs.slide_height)
    
    # suppress unused text placeholder (of index 1, 0 is used for the title)
    sp = shapes[1].element
    sp.getparent().remove(sp)
    
    # GLOBAL results 
    # 1) best team per region
    # create image
    regions_map_img = create_best_team_per_region_img(result_dir, date_of_run, regions_map_file,
                                                      best_teams_per_region)
    
    # init. slide with title
    slide, shapes, title_shape = \
        init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                  "Best team(s) per region (PV prof.)", font_name,
                                  font_size, font_bold, font_italic, text_vertical_align)

    # add image
    add_img_to_slide(slide, regions_map_img, 
                     os.path.join(result_dir, "best_team_per_region_%s.%s" \
                                       % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                          img_format)),
                     ((1-2*left_empty_space/100)*prs.slide_width,
                      (1-2*bottom_empty_space/100)*prs.slide_height \
                                                        -title_shape.height),
                      title_shape.height, left_empty_space/100*prs.slide_width,
                      bottom_empty_space/100*prs.slide_height)
    
    # suppress unused text placeholder (of index 1, 0 is used for the title)
    sp = shapes[1].element
    sp.getparent().remove(sp)
    
    # 2) table with scores per region for all teams (team names in lines, region
    # names in col. and associated score in each cell)
    # TODO put the table in the right placeholder...
    # TODO adapt location and size automatically
    # ALternative for now...
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Scores per team name*region (PV prof.)"
    subtitle.text = "See file: aggreg_per_region_res_run_%s.csv" \
                                       % date_of_run.strftime("%Y-%m-%d_%H%M") 

#    loc_x_table = 600
#    loc_y_table = 500
#    table_width = 2000
#    table_height = 600
#    # init. slide with title
#    slide, shapes, title_shape = \
#        init_img_plus_title_slide(prs, table_slide_layout_idx, 
#                                  "Scores per team name*region (PV prof.)", font_name,
#                                  font_size, font_bold, font_italic, text_vertical_align)
#    # get team and region names based on team_scores
#    team_names = list(team_scores)
#    n_teams = len(team_names)
#    regions = [list(team_scores[team]) for team in team_names]
#    
#    # check that all teams have the same regions scenarios
#    if not check_if_unique_list(regions):
#        print("Different teams does not have the same region names... -> STOP")
#        sys.exit(1)
#    else:
#        regions = regions[0]
#
#    n_regions = len(regions)
    
    # initialize table
#    table = slide.placeholders[1].insert_table(len(team_names)+1, len(regions)+1)
#    table = slide.shapes.add_table(len(team_names)+1, len(regions)+1, loc_x_table,
#                                   loc_y_table, table_width, table_height).table
    # fill each cell by looping over all of them
    # title one
#    table.cell(0, 0).text = "Team/region names"
#    # names of the teams in lines
#    for i_team in range(n_teams):
#        table.cell(i_team+1, 0).text = team_names[i_team]
#    # idem for regions in columns
#    for i_region in range(n_regions):
#        table.cell(0, i_region+1).text = regions[i_region] if not \
#        regions[i_region] in short_region_names else short_region_names[region[i_region]]
#    # loop over (team, region) to fullfill the center cells
#    for i_team in range(n_teams):
#        for i_region in range(n_regions):
#            table.cell(i_team+1, i_region+1).text = \
#                       "%.3f" % team_scores[team_names[i_team]][regions[i_region]]

    # 3) best team France
    # create image
    podium_france_img = create_podium_of_france_img(result_dir, date_of_run, podium_france_file,
                                                    teams_france_classif, type_of_score)
    # init. slide with title
    slide, shapes, title_shape = \
        init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                  "Podium over whole France results", font_name,
                                  font_size, font_bold, font_italic, text_vertical_align)
    
    # add image
    add_img_to_slide(slide, podium_france_img, 
                     os.path.join(result_dir, "podium_france_%s.%s" \
                                       % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                          img_format)),
                     ((1-2*left_empty_space/100)*prs.slide_width,
                      (1-2*bottom_empty_space/100)*prs.slide_height \
                                                        -title_shape.height),
                      title_shape.height, left_empty_space/100*prs.slide_width,
                      bottom_empty_space/100*prs.slide_height)
    
    # suppress unused text placeholder (of index 1, 0 is used for the title)
    sp = shapes[1].element
    sp.getparent().remove(sp)
    
    # improvement slide based on multiple aggreg_results.csv files
    # first check if there is more than one run available
    n_run_dates = len(scores_traj[list(scores_traj)[0]])
    if n_run_dates <= 1:
        print("No slide for score improvement, because unique run available for now")
    else:
        slide, shapes, title_shape = \
            init_img_plus_title_slide(prs, img_slide_layout_idx, 
                                      "All teams score traj. with %s" \
                                      % coord_method, font_name, font_size, font_bold, 
                                      font_italic, text_vertical_align)
    
        # plot and save 
        all_teams_score_traj_file = os.path.join(result_dir,
                                                 "all_teams_score_traj_to_run_%s.%s" \
                                               % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                                  img_format))
    
        plot_all_teams_score_traj(scores_traj, all_teams_score_traj_file.split(".")[0])
        
        # open as Pillow Image
        all_teams_score_traj_img = Image.open(all_teams_score_traj_file)
        
        # add image
        add_img_to_slide(slide, all_teams_score_traj_img, all_teams_score_traj_file,
                         ((1-2*left_empty_space/100)*prs.slide_width,
                          (1-2*bottom_empty_space/100)*prs.slide_height \
                                                            -title_shape.height),
                          title_shape.height, left_empty_space/100*prs.slide_width,
                          bottom_empty_space/100*prs.slide_height)
        
        # suppress unused text placeholder (of index 1, 0 is used for the title)
        sp = shapes[1].element
        sp.getparent().remove(sp)
    

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Annexes"
    subtitle.text = "results for every region"

    for region_coord_dyn_plot in regions:

        # TODO 2021-5-17 exclude teams with bug
        # slide to list these teams

        # 1) 1 slide per team with per iteration total microgrid load
        for i_team in range(n_teams):
            team = team_names[i_team]
            slide, shapes, title_shape = \
                init_img_plus_title_slide(prs, img_slide_layout_idx,
                                          "Team %s load DURING %s \n (PV) region: %s" \
                                          % (team, coord_method, region_coord_dyn_plot),
                                          font_name, font_size, font_bold, font_italic,
                                          text_vertical_align)

            # plot and save
            current_dyn_mg_load_file = os.path.join(result_dir,
                                                    "mg_load_during_dyn_team_%s_%s.%s" \
                                                    % (team, date_of_run.strftime("%Y-%m-%d_%H%M"),
                                                       img_format))
            plot_mg_load_during_coord_method(microgrid_prof, region_coord_dyn_plot,
                                             team, current_dyn_mg_load_file.split(".")[0],
                                             optim_period)

            # open as Pillow Image
            dyn_mg_load_img = Image.open(current_dyn_mg_load_file)

            # add image
            add_img_to_slide(slide, dyn_mg_load_img, current_dyn_mg_load_file,
                             ((1 - 2 * left_empty_space / 100) * prs.slide_width,
                              (1 - 2 * bottom_empty_space / 100) * prs.slide_height \
                              - title_shape.height),
                             title_shape.height, left_empty_space / 100 * prs.slide_width,
                             bottom_empty_space / 100 * prs.slide_height)

            # suppress unused text placeholder (of index 1, 0 is used for the title)
            sp = shapes[1].element
            sp.getparent().remove(sp)

        # 1) 1 slide per team with last iteration per-actor load
        # TODO add possibility to plot best iter
        for i_team in range(n_teams):
            team = team_names[i_team]
            slide, shapes, title_shape = \
                init_img_plus_title_slide(prs, img_slide_layout_idx,
                                          "Team %s ACTORS' load AT THE END of %s \n (PV) region: %s" \
                                          % (team, coord_method, region_coord_dyn_plot),
                                          font_name, font_size, font_bold, font_italic,
                                          text_vertical_align)

            # plot and save
            per_actor_load_file = os.path.join(result_dir,
                                               "per_actor_load_last_iter_team_%s_%s.%s" \
                                               % (team, date_of_run.strftime("%Y-%m-%d_%H%M"),
                                                  img_format))
            plot_per_actor_load_last_iter(load_profiles, pv_prof[region_coord_dyn_plot], region_coord_dyn_plot,
                                          team, per_actor_load_file.split(".")[0],
                                          optim_period)
            # open as Pillow Image
            per_actor_load_img = Image.open(per_actor_load_file)

            # add image
            add_img_to_slide(slide, per_actor_load_img, per_actor_load_file,
                             ((1 - 2 * left_empty_space / 100) * prs.slide_width,
                              (1 - 2 * bottom_empty_space / 100) * prs.slide_height \
                              - title_shape.height),
                             title_shape.height, left_empty_space / 100 * prs.slide_width,
                             bottom_empty_space / 100 * prs.slide_height)

            # suppress unused text placeholder (of index 1, 0 is used for the title)
            sp = shapes[1].element
            sp.getparent().remove(sp)

        # 1 slide with all teams microgrid load at the last iteration
        # TODO add possibility to plot best iter
        slide, shapes, title_shape = \
            init_img_plus_title_slide(prs, img_slide_layout_idx,
                                      "All teams load AT THE END of %s \n (PV) region: %s" \
                                      % (coord_method, region_coord_dyn_plot),
                                      font_name, font_size, font_bold, font_italic,
                                      text_vertical_align)

        # plot and save
        all_teams_mg_load_file = os.path.join(result_dir, "all_teams_mg_load_last_iter_%s.%s" \
                                              % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                                 img_format))
        plot_all_teams_mg_load_last_iter(microgrid_prof, microgrid_pmax, pv_prof[region_coord_dyn_plot],
                                         region_coord_dyn_plot,
                                         all_teams_mg_load_file.split(".")[0],
                                         optim_period)

        # open as Pillow Image
        all_teams_mg_load_img = Image.open(all_teams_mg_load_file)

        # add image
        add_img_to_slide(slide, all_teams_mg_load_img, all_teams_mg_load_file,
                         ((1 - 2 * left_empty_space / 100) * prs.slide_width,
                          (1 - 2 * bottom_empty_space / 100) * prs.slide_height \
                          - title_shape.height),
                         title_shape.height, left_empty_space / 100 * prs.slide_width,
                         bottom_empty_space / 100 * prs.slide_height)

        # suppress unused text placeholder (of index 1, 0 is used for the title)
        sp = shapes[1].element
        sp.getparent().remove(sp)

    # save the pptx JDP presentation
    prs.save(os.path.join(result_dir, "run_summary_%s.pptx" \
                                          % date_of_run.strftime("%Y-%m-%d_%H%M")))

def create_best_team_per_region_img(result_dir, date_of_run: datetime.datetime, regions_map_file: str,
                                    best_teams_per_region: dict) -> Image:
    """
    Create image giving the best team for each of the regions
    
    :param date_of_run: date of the current run 
    :param regions_map_file: full path to the (French) regions map file
    :param best_teams_per_region: dict. with keys the names of the region simulated
    in current run and values the associated list of best team(s)
    """

    # parameters to locate text on the (French) regions map
    # coordinates of the "centers" of the different regions
    regions_center = {"grand_nord": (270,70), "grand_est": (450,170),
                      "grand_rhone": (425,325), "bretagne": (50,150),
                      "grand_ouest": (232,265), "grand_sud_ouest": (230,455),
                      "grande_ardeche": (455, 427), "grand_sud_est": (440,520)}
    regions_textbox_max_width = {"grand_nord": 50, "grand_est": 120,
                                 "grand_rhone": 86, "bretagne": 120,
                                 "grand_ouest": 110, "grand_sud_ouest": 90,
                                 "grande_ardeche": 75, "grand_sud_est": 90}
    
    img_format = "png"
    img_margin_x = 5
    img_margin_y = 5
    textbox_margin_x = 10 # margin box rounding text
    textbox_margin_y = 10
    # TODO logo of given srious gamme (Ponts + 2021?) 
#    logo_relative_size = 10 # in proportion of photo image size
#    logo_margin_x = 5
#    logo_margin_y = 5
#    logo_background_color = ""
#    logo_opacity_level = 0
    
    # for team name style
    font_style = {"name": "Calibri", "size": 15, "bold": True, "italic": False}
    # and associated textbox
    textbox_style = {"outline_color": "black", "outline_width": 2,
                     "outline_style": "solid", "fill_color": "white"}
        
    # open Pillow Image
    regions_map_img = Image.open(regions_map_file)
    draw_regions_map_img = ImageDraw.Draw(regions_map_img)
    
    for region in best_teams_per_region:
        # Join text if multiple teams with same score
        best_teams_txt = "/".join(best_teams_per_region[region])
        # get text width and height
        current_font = ImageFont.truetype("%s%s.ttf" % (font_style["name"].lower(),
                                          "b" if font_style["bold"] else "i" \
                                          if font_style["italic"] else ""),
                                          font_style["size"])
        
        # add linebreak in text to fit with max. textbox width
        best_teams_txt = add_linebreak_to_txt(best_teams_txt, draw_regions_map_img,
                                              current_font, regions_textbox_max_width[region])
        textwidth, textheight = draw_regions_map_img.textsize(best_teams_txt,
                                                              current_font)
        # set text location
        loc_x, loc_y = set_txt_location(regions_center[region], regions_map_img.width,
                                        regions_map_img.height, textwidth, 
                                        textheight, textbox_margin_x+img_margin_x,
                                        textbox_margin_y+img_margin_y)
        
        # textbox rounded the team name(s)
        draw_regions_map_img \
            .rectangle((loc_x-textbox_margin_x, loc_y-textbox_margin_y,
                        loc_x+textwidth+textbox_margin_x,
                        loc_y+textheight+textbox_margin_y),
                        outline=textbox_style["outline_color"],
                        width=textbox_style["outline_width"],
                        fill=textbox_style["fill_color"])

        # Add text
        draw_regions_map_img.text((loc_x, loc_y), best_teams_txt, fill=(0,0,0),
                                  font=current_font)
        
    # save image with best team name(s)     
    regions_map_img \
        .save(os.path.join(result_dir, "best_team_per_region_%s.%s" \
                                       % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                          img_format)))
    
    return regions_map_img

def create_podium_of_france_img(result_dir, date_of_run: datetime.datetime,
                                podium_france_file: str, teams_france_classif: dict,
                                type_of_score: str) -> Image:
    """
    Create image giving the podium best team for each of the regions
    
    :param date_of_run: date of the current run 
    :param podium_france_file: full path to the image of the "podium France" file
    :param teams_france_classif: dict. with keys the different rankings (from 1
    to the number of teams) and values the couple (list of the name(s) of the team(s)),
    score, which may be a number or a vector)
    :param type_of_score: the type of score used for classification (unique 
    number, or a vector)
    :return: returns the Pillow image of the podium
    """

    # parameters to locate text on the (French) regions map
    # coordinates of the "centers" of the different regions
    podium_positions = {1: (705,425), 2: (515,475), 3: (850, 505)}
    podium_textbox_max_width = {1: 120, 2: 100, 3: 110}
    type_of_score_loc = (1220, 500)
    type_of_score_max_width = 314
    
    img_format = "png"
    img_margin_x = 5
    img_margin_y = 5
    textbox_margin_x = 10 # margin box rounding text
    textbox_margin_y = 10
    # TODO logo of given srious gamme (Ponts + 2021?) 
#    logo_relative_size = 10 # in proportion of photo image size
#    logo_margin_x = 5
#    logo_margin_y = 5
#    logo_background_color = ""
#    logo_opacity_level = 0
    
    # for team name style
    font_style = {"name": "Calibri", "size": 20, "bold": True, "italic": False}
    type_of_score_font_style = {"name": "Calibri", "size": 22,
                                "bold": False, "italic": True}
    # and associated textbox
    textbox_style = {"outline_color": "black", "outline_width": 2,
                     "outline_style": "solid", "fill_color": "white"}
        
    # open Pillow Image
    podium_france_img = Image.open(podium_france_file)
    draw_podium_france_img = ImageDraw.Draw(podium_france_img)

    # add team names on the podium
    rankings = [r for r in teams_france_classif.keys() if r in [1, 2, 3]]
    for ranking in rankings:
        # Join text if multiple teams with same score
        given_ranking_teams_txt = "%.2f: %s" % (teams_france_classif[ranking][0],
                                                teams_france_classif[ranking][1])
        # get text width and height
        current_font = ImageFont.truetype("%s%s.ttf" % (font_style["name"].lower(),
                                          "b" if font_style["bold"] else "i" \
                                          if font_style["italic"] else ""),
                                          font_style["size"])
        
        # add linebreak in text to fit with max. textbox width
        given_ranking_teams_txt = \
            add_linebreak_to_txt(given_ranking_teams_txt, draw_podium_france_img,
                                 current_font, podium_textbox_max_width[ranking])
        textwidth, textheight = draw_podium_france_img.textsize(given_ranking_teams_txt,
                                                                current_font)
        # set text location
        loc_x, loc_y = set_txt_location(podium_positions[ranking], podium_france_img.width,
                                        podium_france_img.height, textwidth, 
                                        textheight, textbox_margin_x+img_margin_x,
                                        textbox_margin_y+img_margin_y)
        
        # textbox rounded the team name(s)
        draw_podium_france_img \
            .rectangle((loc_x-textbox_margin_x, loc_y-textbox_margin_y,
                        loc_x+textwidth+textbox_margin_x,
                        loc_y+textheight+textbox_margin_y),
                        outline=textbox_style["outline_color"],
                        width=textbox_style["outline_width"],
                        fill=textbox_style["fill_color"])

        # Add text
        draw_podium_france_img.text((loc_x, loc_y), given_ranking_teams_txt,
                                    fill=(0,0,0), font=current_font)

    # add the type of score legend
    # get text width and height
    current_font = ImageFont.truetype("%s%s.ttf" % (type_of_score_font_style["name"].lower(),
                                      "b" if type_of_score_font_style["bold"] else "i" \
                                      if type_of_score_font_style["italic"] else ""),
                                      type_of_score_font_style["size"])

    # add linebreak in text to fit with max. textbox width
    type_of_score_txt = "SCORE (to be %s): %s" % (type_of_score[0], type_of_score[1])
    type_of_score_txt = \
        add_linebreak_to_txt(type_of_score_txt, draw_podium_france_img,
                             current_font, type_of_score_max_width)
    textwidth, textheight = draw_podium_france_img.textsize(type_of_score_txt,
                                                            current_font)
    # set text location
    loc_x, loc_y = set_txt_location(type_of_score_loc, podium_france_img.width,
                                    podium_france_img.height, textwidth, 
                                    textheight, textbox_margin_x+img_margin_x,
                                    textbox_margin_y+img_margin_y)

    # textbox rounded the type of score legend
    draw_podium_france_img \
        .rectangle((loc_x-textbox_margin_x, loc_y-textbox_margin_y,
                    loc_x+textwidth+textbox_margin_x,
                    loc_y+textheight+textbox_margin_y),
                    outline=textbox_style["outline_color"],
                    width=textbox_style["outline_width"],
                    fill=textbox_style["fill_color"])

    # Add text
    draw_podium_france_img.text((loc_x, loc_y), type_of_score_txt,
                                fill=(0,0,0), font=current_font)
        
    # save image with best team name(s)     
    podium_france_img \
        .save(os.path.join(result_dir, "podium_france_%s.%s" \
                                       % (date_of_run.strftime("%Y-%m-%d_%H%M"),
                                          img_format)))
    
    return podium_france_img

def set_txt_location(txt_center_ideal_loc: tuple, img_width: (float, int), 
                     img_height: (float, int), textwidth: (float, int),
                     textheight: (float, int), margin_x: (float, int),
                     margin_y: (float, int)) -> tuple:
    """
    Set text location
    
    :param txt_center_ideal_loc: (x,y) ideal location for the center of the 
    textbox
    :param img_width: image width on which the text will be written
    :param img_height: idem height
    :param textwidth: text width
    :param textheight: text height
    :param margin_x: x-axis margin to be taken between text and the image limits
    :param margin_x: idem for y-axis    """
    
    # x location
    if textwidth + 2*margin_x > img_width:
        print("TextWIDTH pb: the text to be added is bigger than image!")
        loc_x = txt_center_ideal_loc[0] - textwidth/2
    else:
        # check if text width is in conflict with left/right limits of the image
        loc_x = txt_center_ideal_loc[0] - textwidth/2 if \
            (txt_center_ideal_loc[0] - textwidth/2 - margin_x >= 0 and \
                     txt_center_ideal_loc[0] + textwidth/2 + margin_x <= img_width) \
             else margin_x if txt_center_ideal_loc[0] - textwidth/2 - margin_x < 0 \
                             else img_width - textwidth - margin_x

    # y location
    if textheight + 2*margin_y> img_height:
        print("TextHEIGHT pb: the text to be added is bigger than image!")
        loc_y = txt_center_ideal_loc[1] - textheight/2
    else:
        # check if text height conflicting with top/bottom img limits
        loc_y = txt_center_ideal_loc[1] - textheight/2 if \
            (txt_center_ideal_loc[1] - textheight/2 - margin_y >= 0 and \
                     txt_center_ideal_loc[1] + textheight/2 + margin_y <= img_height) \
             else margin_y if txt_center_ideal_loc[1] - textheight/2 - margin_y < 0 \
                          else img_height - textheight - margin_y

    return loc_x, loc_y

def add_linebreak_to_txt(my_txt: str, img_draw: ImageDraw, text_font: ImageFont,
                         max_width: int) -> str:
    """
    Add linebreak in text to fit with a max. width limit
    """
    
    textwidth, textheight = img_draw.textsize(my_txt, text_font)
    
    if textwidth <= max_width:
        return my_txt
    else:
        # get number of characters which correspond to a line
        line_len = int(np.floor(max_width/textwidth*len(my_txt)))
        
        my_txt_with_linebreaks = ""
        i = 0
        while i < len(my_txt):
            current_line = my_txt[i:i+line_len].lstrip()
            if len(current_line) < line_len:
                delta_fullfill = line_len - len(current_line)
                current_line += my_txt[i+line_len:i+line_len+delta_fullfill]
            else:
                delta_fullfill = 0
                
            my_txt_with_linebreaks += "%s\n" % current_line
            i += line_len + delta_fullfill
        
        # suppress last linebreak        
        return my_txt_with_linebreaks[:-1]

def add_img_to_slide(slide: pptx.slide.Slide, img_to_add: Image, img_file: str,
                     img_box: tuple, title_height: float, left_space: float,
                     bottom_space: float):
    """
    Add an image to a simple slide with a title and an image
    
    :param slide: slide in which image has to be added 
    :param img_to_add: Pillow Image to be added
    :param img_file: full path to the file of this image
    :param img_box: image box dimensions (width, height)
    :param title_height: height of the title above the image in slide
    :param left_space: space to the left of the img_box in slide
    :param bottom_space: space at the bottom of the img_box (and above) in slide
    """
    
    img_resize = resize_img_in_box(img_to_add, img_box)
    
    # set location based on img new dimensions
    left_img = left_space + (img_box[0]-img_resize[0])/2
    top_img = title_height + bottom_space + (img_box[1]-img_resize[1])/2
    # add this image
    slide.shapes.add_picture(img_file, left_img, top_img, height=img_resize[1]) 

def init_img_plus_title_slide(prs: Presentation, layout_idx: int, title_text: str,
                              font_name: str, font_size: int, font_bold: bool, 
                              font_italic: bool, text_vertical_align: str):
    """
    Initialize ppt slide with title and an image
    """
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = title_text
    # set title style
    set_text_style(title_shape.text_frame.paragraphs[0], font_name, font_size,
                   font_bold, font_italic, text_vertical_align)
    
    return slide, shapes, title_shape

def set_text_style(text_shape, font_name: str="Calibri", font_size: int=18,
                   font_bold: bool=True, font_italic: bool=False, 
                   text_vertical_align: str="center"):
    """
    Set text style of a given text shape
    """
    font = text_shape.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = font_bold
    font.italic = font_italic  # cause value to be inherited from theme
    if text_vertical_align == "top":
        text_shape.vertical_anchor = MSO_ANCHOR.TOP
    elif text_vertical_align == "middle":
        text_shape.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif text_vertical_align == "bottom":
        text_shape.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
# ERROR
        print("Unknown text vertical alignment %s -> STOP" % text_vertical_align)
        sys.exit(1)

def resize_img_in_box(img: Image, img_box_size: tuple) -> tuple:
    """
    Resize image based on the size of an associated box
    
    :param img: Pillow image
    :param img_box_size: (width, height) of the associated box
    :return: returns the tuple (width, height) with updated image dimensions
    """
    
    # sizing operation limited by height
    if img_box_size[0] / img.width >= img_box_size[1] / img.height:
        return (img.width*img_box_size[1]/img.height, img_box_size[1])
    else:
        return (img_box_size[0], img.height*img_box_size[0]/img.width)

def create_current_run_dir(current_dir: str, date_of_run: datetime.datetime):
    """
    Create a directory where all results of current run will be saved
    
    :param current_dir: directory in which photos are stored
    :param date_of_run: date of current run
    :return: returns the full path of current run directory
    """
    
    current_run_dir = os.path.join(current_dir,
                                   "run_%s" % date_of_run.strftime("%Y-%m-%d_%H%M"))
    
    if os.path.exists(current_run_dir):
        print("Result dir. for run %s already exists -> it is emptied before processing current outputs" \
              % date_of_run.strftime("%Y-%m-%d_%H%M"))
        all_files_tb_suppr = [os.path.join(current_run_dir, elt_file) \
                            for elt_file in os.listdir(current_run_dir) \
                            if (elt_file.endswith(".png") or elt_file.endswith(".csv"))]
        
        for elt_file in all_files_tb_suppr:
            os.remove(elt_file)
    
    # create backup photos dir
    else:
        os.mkdir(current_run_dir)
    
    return current_run_dir

if __name__ =="__main__":
    current_dir = os.getcwd()
    date_of_run = datetime.datetime.now()
    idx_run = 1
    
    # create directory for current run
    result_dir = create_current_run_dir(current_dir, date_of_run)
    
    # TODO unify team names (etc.) in the different tests
    
#    best_teams_per_region = {"grand_nord": ["team 1"], "grand_est": ["big team"], 
#                             "grand_rhone": ["qui ne saute pas"], 
#                             "bretagne": ["bonnets rouges et blanc bonnet"], 
#                             "grand_ouest": ["a la masse"], 
#                             "grand_sud_ouest": ["la garonne est viola"],
#                             "grande_ardeche": ["ca canyone pas mal"], 
#                             "grand_sud_est": ["le soleil c'est la vie"]}
    regions_map_file = os.path.join(current_dir, "images", "pv_regions_no_names.png")
#    # Test creation of best team per region image
#    create_best_team_per_region_img(date_of_run, regions_map_file,
#                                    best_teams_per_region)
    
    podium_france_file = os.path.join(current_dir, "images", "podium_france_v2.png")
    teams_france_classif = {1: (["best of the best OM forever"], 800),
                            2: (["not so bad", "same perf than you"], 990),
                            3: (["happy to be here"], 1000)}
    autonomy_weight = 3.1415
    type_of_score = ("MIN.", 
                     "Weighted (bill + contract. pmax) + %.2f * autonomy score" \
                        % autonomy_weight) 
#    # Test creation of France podium image
#    create_podium_of_france_img(date_of_run, podium_france_file, 
#                                teams_france_classif, type_of_score)
    
    team_scores = {"team_1": {"grand_nord": 12, "grand_est": 57, "grand_rhone": 69, 
                              "bretagne": 35, "grand_ouest": 47, "grand_sud_ouest": 33,
                              "grande_ardeche": 72, "grand_sud_est": 22},
                   "team_2": {"grand_nord": 15, "grand_est": 52, "grand_rhone": 61,
                              "bretagne": 32, "grand_ouest": 41, "grand_sud_ouest": 31,
                              "grande_ardeche": 71, "grand_sud_est": 22},
                   "team_3": {"grand_nord": 9, "grand_est": 57, "grand_rhone": 75, 
                              "bretagne": 30, "grand_ouest": 40, "grand_sud_ouest": 29,
                              "grande_ardeche": 56, "grand_sud_est": 18}
                  }
    
    delta_t_s = 1800 # time-slot duration (s)
    start_optim_period = datetime.datetime(2018,1,1)
    optim_period = pd.date_range(start=start_optim_period, 
                                 end=start_optim_period+timedelta(hours=24),
                                 freq="%is" % delta_t_s)[:-1]
    coord_method = "price-coord. dyn."
    # Collective metrics calculation
    n_ts = 48
    load_profiles = {1: 
                       {1: 
                          {"grand_nord": 
                                {1: 
                                   {"team_1": 
                                            {1: 
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             2:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             3:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)}
                                             },
                                    "team_2":
                                            {1: 
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             2:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             3:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             4:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)}
                                             },
                                    "team_3":
                                            {1: 
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             2:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             3:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)},
                                             4:
                                               {"solar_farm": np.random.rand(n_ts),
                                                "industrial_consumer": np.random.rand(n_ts),
                                                "charging_station": np.random.rand(n_ts),
                                                "data_center": np.random.rand(n_ts)}
                                             }
                                   }
                               }
                         }
                      }
                     }

    purchase_price = 0.10 + 0.1 * np.random.rand(n_ts)
    sale_price = 0.05 + 0.1 * np.random.rand(n_ts)
    delta_t_s = 1800
    contracted_p_tariffs = {6: 123.6, 9: 151.32, 12: 177.24, 15: 201.36,
                            18: 223.68, 24: 274.68, 30: 299.52, 36: 337.56}

    from calc_output_metrics import calc_per_actor_bills, calc_microgrid_collective_metrics, \
       calc_cost_autonomy_tradeoff_last_iter, get_best_team_per_region, \
           save_all_metrics_to_csv, save_per_region_score_to_csv, get_improvement_traj

    # calculate per-actor bill
    per_actor_bills = calc_per_actor_bills(load_profiles, purchase_price,
                                           sale_price, delta_t_s)

    # and microgrid collective metrics
    microgrid_prof, microgrid_pmax, collective_metrics = \
            calc_microgrid_collective_metrics(load_profiles, contracted_p_tariffs, 
                                              delta_t_s)

    # and finally cost-autonomy tradeoff
    cost_autonomy_tradeoff = \
        calc_cost_autonomy_tradeoff_last_iter(per_actor_bills, collective_metrics)

    # Get best team per region
    coll_metrics_weights = {"pmax_cost": 1/365, "autonomy_score": 1,
                            "mg_transfo_aging": 0, "n_disj": 0}
    team_scores, best_teams_per_region, coll_metrics_names = \
            get_best_team_per_region(per_actor_bills, collective_metrics,
                                     coll_metrics_weights)
    
    # save detailed results to a .csv file
    metrics_not_saved = ["mg_transfo_aging", "n_disj"]
    save_all_metrics_to_csv(per_actor_bills, collective_metrics, coll_metrics_names,
                            coll_metrics_weights, metrics_not_saved, result_dir,
                            date_of_run)
    
    # and aggreg. per region .csv file
    save_per_region_score_to_csv(team_scores, result_dir, date_of_run)
    
    # get "improvement trajectory"
    list_of_run_dates = [datetime.datetime.strptime(elt[4:], "%Y-%m-%d_%H%M") \
                         for elt in os.listdir(current_dir) \
                         if (os.path.isdir(elt) and elt.startswith("run_"))]
    scores_traj = get_improvement_traj(current_dir, list_of_run_dates,
                                       list(team_scores))

    pv_prof = 5*np.random.rand(n_ts)

    create_summary_of_run_ppt(result_dir, date_of_run, idx_run, optim_period,
                              coord_method, regions_map_file, pv_prof, load_profiles,
                              microgrid_prof, microgrid_pmax, cost_autonomy_tradeoff,
                              team_scores, best_teams_per_region, podium_france_file,
                              teams_france_classif, type_of_score, scores_traj)
    